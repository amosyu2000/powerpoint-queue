import os
import pptx

TITLE_SLIDE_LAYOUT = 0

# Set path to the path of the project
directory_path = os.path.abspath(os.path.dirname(__file__))
os.chdir(directory_path)
os.chdir("../")

# Open the pptx file
presentation = pptx.Presentation('res/template.pptx')
slide_layout = presentation.slide_layouts[TITLE_SLIDE_LAYOUT]

# Create a new slide
slide = presentation.slides.add_slide(slide_layout)

# Populate each placeholder in the slide with its corresponding idx value
for shape in slide.placeholders:
	idx = shape.placeholder_format.idx
	print('idx: %d, name: %s' % (idx, shape.name))
	shape.text = str(idx)

# Save the presentation in the tests directory
presentation.save('tests/test.pptx')
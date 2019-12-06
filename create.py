import os
import pptx

# ==========
# CONSTANTS
# ==========

TITLE_SLIDE_LAYOUT = 0
QUEUE_POS = 10
SLIDE_NUMBER = 15
DEST = 'C:/Users/Dell/Downloads'


def run():

	# Set path to the path of the project
	directory_path = os.path.abspath(os.path.dirname(__file__))
	os.chdir(directory_path)

	# Open file containing queue
	infile = open('res/q.txt', 'r')

	# Parse queue into list by splitting at the new-line characters
	queue = infile.read().split('\n')

	# Open the pptx file
	presentation = pptx.Presentation('res/template.pptx')
	slide_layout = presentation.slide_layouts[TITLE_SLIDE_LAYOUT]

	# For each name in the queue
	for i in range(len(queue)):

		# Create a new slide
		slide = presentation.slides.add_slide(slide_layout)

		# Put the slide number
		slide.placeholders[SLIDE_NUMBER].text = str(i+1) + ' of ' + str(len(queue))

		# Put the "Currently Serving" name and the three "Up Next" names
		for pos in range(4):
			try:
				slide.placeholders[QUEUE_POS + pos].text = queue[i+pos]
			# Put empty text if there is no more names in the queue
			except IndexError:
				slide.placeholders[QUEUE_POS + pos].text = ' '

		# If there are more people in the list, add an ellipsis
		try:
			queue[i+4]
			slide.placeholders[QUEUE_POS + 4].text = '...'
		except IndexError:
			slide.placeholders[QUEUE_POS + 4].text = ' '

	# Save the presentation in the user-specified destination
	presentation.save(DEST+'/template.pptx')